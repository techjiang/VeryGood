"""pptx 抽取：逐页形状文本、表格、备注页、超链接。

iter_shapes 同时供 defect_check.py 的几何类检查复用（出界/重叠需要活对象，只有文本
抽取走本模块，保证解析口径一致）。
"""
from .common import merge_links, ooxml_external_links, text_urls

GROUP_SHAPE_TYPE = 6  # MSO_SHAPE_TYPE.GROUP


def iter_shapes(shapes, depth=0):
    """展开组合形状（限深 3，防异常嵌套）。"""
    for sh in shapes:
        yield sh
        if depth < 3:
            try:
                if sh.shape_type == GROUP_SHAPE_TYPE:
                    for inner in iter_shapes(sh.shapes, depth + 1):
                        yield inner
            except Exception:
                pass


def shape_text(sh):
    try:
        if sh.has_text_frame:
            return (sh.text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def extract(path):
    units, outline, notes = [], [], []
    try:
        from pptx import Presentation
    except Exception:
        return {
            "units": [],
            "links": merge_links(ooxml_external_links(path)),
            "outline": [],
            "meta": {},
            "notes": ["python-pptx 不可用，pptx 正文未抽取（pip install python-pptx）"],
        }
    try:
        prs = Presentation(path)
    except Exception as e:
        return {"units": [], "links": [], "outline": [], "meta": {}, "notes": [f"pptx 打开失败: {e}"]}

    slides = list(prs.slides)
    text_links = []
    for si, slide in enumerate(slides, 1):
        first = None
        for sh in iter_shapes(slide.shapes):
            t = shape_text(sh)
            if t:
                units.append({"anchor": f"pptx:s{si}", "text": t})
                text_links += text_urls(t, f"第{si}页")
                if first is None:
                    first = t
            try:
                if sh.has_table:
                    for ri, row in enumerate(sh.table.rows, 1):
                        cells = [(c.text or "").replace("\n", " ").strip() for c in row.cells]
                        if any(cells):
                            line = "\t".join(cells)
                            units.append({"anchor": f"pptx:s{si}t r{ri}".replace(" ", ""), "text": line})
                            text_links += text_urls(line, f"第{si}页表格")
            except Exception:
                pass
        try:
            if slide.has_notes_slide:
                nt = (slide.notes_slide.notes_text_frame.text or "").strip()
                if nt:
                    units.append({"anchor": f"pptx:s{si}n", "text": nt})
                    text_links += text_urls(nt, f"第{si}页备注")
        except Exception:
            pass
        outline.append({"anchor": f"pptx:s{si}", "label": (first or "(无文字)")[:80], "level": 1})

    return {
        "units": units,
        "links": merge_links(ooxml_external_links(path), text_links),
        "outline": outline,
        "meta": {"slides": len(slides)},
        "notes": notes,
    }
