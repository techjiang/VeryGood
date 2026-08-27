#!/usr/bin/env python3
"""交付产物的结构 / 内容确定性校验（只读，不修改任何文件）。

用法:
    python3 defect_check.py <file>
    python3 defect_check.py --doctor            # 依赖自检：一条命令输出各依赖可用性
    python3 defect_check.py --doctor --install   # 顺手把缺失依赖按国内源优先装上

向 stdout 输出单个 JSON 对象。本脚本永不抛栈：任何失败都收敛进 JSON，
调用方（校验子代理）可无条件解析结果。

checks[].status 三档语义（调用方据此判定，不可混用）:
    fail = 确定性缺陷（打不开/格式损坏/占位符残留/空片/公式错误值/引用资源缺失/
           html 内联 JS 语法错误）
           → 证据充分，对应 must 判不满足
    warn = 启发式参考信号（字号偏小/文本框重叠/空章节/无文本页/引用可疑/
           图表容器与引库约定/非 utf-8 源码）
           → 记入证据，单独不足以判 must 不满足，需结合内容判断
    skip = 未覆盖（依赖缺失或数据不足）→ 如实标注，既不算通过也不算失败
    pass = 该项已实测且无问题


JSON 结构:
    {
      "file": str, "ext": str,
      "openable": bool,        # 能否被对应库正常打开
      "format_valid": bool,    # 无 fail 级缺陷
      "checks": [ {"name": str, "status": "pass|fail|warn|skip", "detail": str} ],
      ...按类型附加 page_count / slides / sheets / paragraphs / text_chars
    }
"""
import argparse
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
# 解析工具函数与 content_extract.py 共用同一套实现（extractors 包），避免两处口径漂移。
# defect_check.py 仍各自跑一次遍历——它的检查需要活对象（单元格/形状/页面），不是纯文本。
from extractors.docx import heading_level  # noqa: E402
from extractors.pptx import iter_shapes, shape_text  # noqa: E402

try:  # 依赖自举（国内源优先）。import 即把缓存 pylibs 挂进 sys.path，装过的库才 import 得到
    import deps  # noqa: E402
except Exception:  # pragma: no cover - 缺文件时退化为"只用系统已装的"
    deps = None

OFFICE_ERRORS = {"#REF!", "#DIV/0!", "#VALUE!", "#NAME?", "#N/A", "#NULL!", "#NUM!"}

# 占位符残留：交付物里出现即为确定性缺陷（模板默认文字 / 未替换变量 / 待办标记）。
# 命中片段会带上下文回传，便于调用方核实是否误报（如正文本就在讨论 TODO 机制）。
PLACEHOLDER_RES = [
    ("TODO/FIXME", re.compile(r"\b(?:TODO|FIXME|TBD|TBC)\b")),
    ("待补充", re.compile(r"待补充|待填写|待完善|待确认|待定|此处填写|占位符")),
    ("lorem ipsum", re.compile(r"lorem\s+ipsum", re.I)),
    ("XXX", re.compile(r"[xX]{3,}")),
    ("[insert...]", re.compile(r"\[insert[^\]]{0,40}\]", re.I)),
    ("未替换模板变量", re.compile(r"\{\{[^}\n]{1,40}\}\}")),
    ("版式默认文字", re.compile(r"单击此处添加|点击此处添加|单击此处编辑|Click to (?:add|edit)", re.I)),
]

MAX_HITS = 20          # 单项缺陷最多回传条数，避免输出爆炸
MAX_CELLS = 200000     # xlsx 单次扫描单元格上限，防超大表拖死
TINY_FONT_PT = 10.0    # 正文字号下限（pptx，低于此值记 warn）
_TMP_ROOT = Path.cwd() / ".tmp"


def _result(file, ext, openable, format_valid, checks, **extra):
    out = {"file": file, "ext": ext, "openable": openable, "format_valid": format_valid, "checks": checks}
    out.update(extra)
    return out


def _c(name, status, detail=""):
    return {"name": name, "status": status, "detail": detail}


def _has_fail(checks):
    return any(c["status"] == "fail" for c in checks)


def _snip(text, m, width=30):
    """命中片段带上下文，供调用方核实。"""
    s = max(0, m.start() - width)
    return text[s : m.end() + width].replace("\n", " ").strip()


def _scan_placeholders(items):
    """items 为 [(位置标签, 文本), ...]；返回命中描述列表。"""
    hits = []
    for where, text in items:
        if not text:
            continue
        for label, rx in PLACEHOLDER_RES:
            m = rx.search(text)
            if m:
                hits.append(f"{where}: [{label}] {_snip(text, m)}")
                if len(hits) >= MAX_HITS:
                    return hits
    return hits


def _placeholder_check(items):
    hits = _scan_placeholders(items)
    if hits:
        return _c("placeholder", "fail", "; ".join(hits))
    return _c("placeholder", "pass")


# ── soffice 调用（沙箱内需要 browser.py 的 AF_UNIX shim，故复用其环境构造） ──
def _soffice_env():
    try:
        sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
        from browser import soffice_env  # 复用沙箱 LD_PRELOAD shim

        return soffice_env()
    except Exception:
        env = os.environ.copy()
        env["SAL_USE_VCLPLUGIN"] = "svp"
        return env


# LibreOffice 载入 OOXML 时默认不重算公式（OOXMLRecalcMode=ask，headless 下等于 never），
# 于是文件里缓存的错误值/陈旧值都会被原样带过。用一次性 user profile 把重算模式设为
# always(0)，转换时强制重算，再读缓存值才能真正发现 #REF! 这类计算错误。
_LO_RECALC_XCU = """<?xml version="1.0" encoding="UTF-8"?>
<oor:items xmlns:oor="http://openoffice.org/2001/registry" xmlns:xs="http://www.w3.org/2001/XMLSchema">
 <item oor:path="/org.openoffice.Office.Calc/Formula/Load"><prop oor:name="OOXMLRecalcMode" oor:op="fuse"><value>0</value></prop></item>
 <item oor:path="/org.openoffice.Office.Calc/Formula/Load"><prop oor:name="ODFRecalcMode" oor:op="fuse"><value>0</value></prop></item>
</oor:items>
"""


def _work_dir(path, tag):
    """按绝对路径 hash 建隔离目录：不同目录下的同名产物不会互相覆盖。"""
    h = hashlib.sha1(os.path.abspath(path).encode()).hexdigest()[:8]
    stem = re.sub(r"[^\w.-]+", "_", Path(path).stem) or "artifact"
    d = _TMP_ROOT / f"dv-{tag}-{stem}-{h}"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _recalc_xlsx(path):
    """soffice 强制重算后另存一份副本，返回副本路径；不可用时返回 (None, 原因)。"""
    if not shutil.which("soffice"):
        return None, "soffice unavailable"
    try:
        outdir = _work_dir(path, "recalc")
        profile = outdir / "lo-profile"
        (profile / "user").mkdir(parents=True, exist_ok=True)
        (profile / "user" / "registrymodifications.xcu").write_text(_LO_RECALC_XCU, encoding="utf-8")
        r = subprocess.run(
            [
                "soffice",
                f"-env:UserInstallation={profile.as_uri()}",
                "--headless",
                "--norestore",
                "--convert-to",
                "xlsx",
                "--outdir",
                str(outdir),
                str(path),
            ],
            capture_output=True,
            text=True,
            timeout=180,
            env=_soffice_env(),
        )
        out = outdir / f"{Path(path).stem}.xlsx"
        if out.exists():
            return str(out), ""
        return None, f"soffice failed: {(r.stderr or r.stdout).strip()[:160]}"
    except Exception as e:
        return None, f"recalc error: {str(e)[:160]}"


def check_office(path, ext):
    checks = []
    try:
        with zipfile.ZipFile(path) as zf:
            bad = zf.testzip()
            if bad is not None:
                checks.append(_c("zip_integrity", "fail", f"corrupt entry: {bad}"))
                return _result(path, ext, False, False, checks)
            checks.append(_c("zip_integrity", "pass"))
            names = set(zf.namelist())
    except Exception as e:
        checks.append(_c("zip_integrity", "fail", str(e)))
        return _result(path, ext, False, False, checks)

    if "[Content_Types].xml" not in names:
        checks.append(_c("ooxml_marker", "fail", "missing [Content_Types].xml"))
        return _result(path, ext, False, False, checks)
    checks.append(_c("ooxml_marker", "pass"))

    if ext in ("xlsx", "xlsm"):
        return _check_xlsx(path, ext, checks)
    if ext == "docx":
        return _check_docx(path, ext, checks)
    if ext == "pptx":
        return _check_pptx(path, ext, checks)
    return _result(path, ext, True, True, checks)


def _scan_xlsx_raw(wb):
    """一遍扫过 data_only=False 的工作簿：收字面错误文本、公式单元格坐标、文本片段。"""
    hits, formulas, texts = [], {}, []
    scanned = 0
    for ws in wb.worksheets:
        coords, buf = set(), []
        for row in ws.iter_rows():
            for cell in row:
                scanned += 1
                if scanned > MAX_CELLS:
                    break
                v = cell.value
                if not isinstance(v, str):
                    continue
                s = v.strip()
                if s.startswith("="):
                    coords.add(cell.coordinate)
                    continue
                if s in OFFICE_ERRORS and len(hits) < MAX_HITS:
                    hits.append(f"{ws.title}!{cell.coordinate}={s}")
                if len(buf) < 2000:
                    buf.append(v)
        formulas[ws.title] = coords
        texts.append((f"sheet[{ws.title}]", "\n".join(buf)))
    return hits, formulas, texts


def _scan_xlsx_cached(wb, formulas):
    """扫 data_only=True 工作簿里公式单元格的缓存值。
    返回 (错误值命中, 是否至少有一个公式单元格带缓存值)。"""
    hits, seen = [], False
    for ws in wb.worksheets:
        coords = formulas.get(ws.title) or set()
        if not coords:
            continue
        for coord in coords:
            try:
                v = ws[coord].value
            except Exception:
                continue
            if v is None:
                continue
            seen = True
            if isinstance(v, str) and v.strip() in OFFICE_ERRORS and len(hits) < MAX_HITS:
                hits.append(f"{ws.title}!{coord}={v.strip()}")
    return hits, seen


def _chart_labels_check(findings, n_charts):
    """图表数据标注：成分是否重复图例/类目轴 + 相邻标注是否压盖。

    重复与压盖同时成立才记 fail —— 单看"成分未显式关闭"可能是刻意标注，单看压盖是不建模
    自动换行的宽度上界；两者叠加时图表必然既冗余又读不出数值，证据才算充分。
    """
    dup = [f for f in findings if f["redundant"]]
    lap = [f for f in findings if f["overlaps"]]
    if not dup and not lap:
        return _c("chart_labels", "pass", f"{n_charts} 个图表的数据标注无冗余成分、无压盖")
    lines = []
    for f in findings:
        if not f["redundant"] and not f["overlaps"]:
            continue
        bits = [f"{f['where']}: {f['series']}系列×{f['points']}点={f['labels']} 个标注"]
        if f["redundant"]:
            bits.append("标注含 " + "、".join(f["redundant"]) + "，与图例/类目轴内容重复")
        if f["overlaps"]:
            sample = f"（{f['hits'][0]}）" if f["hits"] else ""
            bits.append(f"估算 {f['overlaps']} 处横向压盖{sample}")
        lines.append("，".join(bits))
        if len(lines) >= MAX_HITS:
            break
    detail = "; ".join(lines)
    status = "fail" if [f for f in dup if f["overlaps"]] else "warn"
    if dup:
        detail += "。建议数据标注只保留数值（showVal），显式关闭 showSerName/showCatName/showLegendKey"
    return _c("chart_labels", status, detail)


def _check_xlsx(path, ext, checks):
    try:
        import openpyxl
    except Exception:
        checks.append(_c("library_open", "skip", "openpyxl unavailable"))
        checks.append(_c("formula_error", "skip", "openpyxl unavailable, 公式错误未校验"))
        return _result(path, ext, True, True, checks, formula_check_mode="unavailable")
    try:
        wb = openpyxl.load_workbook(path, data_only=False)
    except Exception as e:
        checks.append(_c("library_open", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    sheets = wb.sheetnames
    checks.append(_c("library_open", "pass", f"{len(sheets)} sheet(s)"))

    # ① 字面错误文本（写死在单元格里的 #REF! 等）+ 公式单元格坐标 + 文本片段
    literal, formulas, texts = _scan_xlsx_raw(wb)
    n_formula = sum(len(v) for v in formulas.values())

    # ② 计算错误：soffice 强制重算 → 读公式单元格缓存值。重算不可用时退到直读原文件缓存值；
    #    若公式单元格一个缓存值都没有（脚本生成的 xlsx 通常不写缓存）→ 如实标 skip，绝不静默通过。
    mode, computed, note = "unavailable", [], ""
    if n_formula == 0:
        mode = "no_formula"
    else:
        recalc_path, why = _recalc_xlsx(path)
        if recalc_path:
            try:
                wb2 = openpyxl.load_workbook(recalc_path, data_only=True)
                computed, seen = _scan_xlsx_cached(wb2, formulas)
                mode = "recalc" if seen else "unavailable"
                if not seen:
                    note = "重算后仍无缓存值"
            except Exception as e:
                note = f"recalc file unreadable: {str(e)[:100]}"
        else:
            note = why
        if mode not in ("recalc",):
            try:
                wb3 = openpyxl.load_workbook(path, data_only=True)
                computed, seen = _scan_xlsx_cached(wb3, formulas)
                if seen:
                    mode = "cached"
            except Exception as e:
                note = (note + "; " if note else "") + f"cached read failed: {str(e)[:100]}"

    hits = sorted(set(literal) | set(computed))
    if hits:
        checks.append(_c("formula_error", "fail", "; ".join(hits[:MAX_HITS])))
    elif mode == "no_formula":
        checks.append(_c("formula_error", "pass", "无公式单元格"))
    elif mode == "recalc":
        checks.append(_c("formula_error", "pass", f"soffice 重算 {n_formula} 个公式后无错误值"))
    elif mode == "cached":
        checks.append(_c("formula_error", "pass", f"基于文件内缓存值判定（未重算：{note or 'recalc unavailable'}）"))
    else:
        checks.append(
            _c(
                "formula_error",
                "skip",
                f"{n_formula} 个公式单元格无缓存值且无法重算，计算错误未校验（{note or 'no cached values'}）",
            )
        )

    checks.append(_placeholder_check(texts))

    # ③ 图表数据标注：成分冗余（重复图例/类目轴）+ 相邻标注压盖
    n_charts = 0
    try:
        from checkers.xlsx_charts import audit as audit_chart_labels

        findings, n_charts = audit_chart_labels(wb)
    except Exception as e:
        checks.append(_c("chart_labels", "skip", f"图表标注未校验: {str(e)[:120]}"))
    else:
        if n_charts:
            checks.append(_chart_labels_check(findings, n_charts))

    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks, sheets=sheets, charts=n_charts, formula_check_mode=mode)


def _docx_texts(d):
    items, total = [], 0
    for i, p in enumerate(d.paragraphs):
        t = p.text or ""
        total += len(t.strip())
        if t.strip():
            items.append((f"段落{i + 1}", t))
    for ti, tbl in enumerate(d.tables):
        for ri, row in enumerate(tbl.rows):
            for ci, cell in enumerate(row.cells):
                t = cell.text or ""
                total += len(t.strip())
                if t.strip():
                    items.append((f"表{ti + 1}[{ri + 1},{ci + 1}]", t))
    return items, total


def _docx_images(d):
    """body 里的图片数，口径与 `docx_pages.count_images` 一致。它决定 `non_empty` 判 fail
    还是 pass —— 口径分家就会把图片型文档判成空文档。"""
    try:
        from checkers import docx_pages

        return docx_pages.count_images(d)
    except Exception:
        try:
            return len(d.inline_shapes)
        except Exception:
            return 0


def _heading_level(p):
    return heading_level(p)


def _docx_body_seq(d):
    """按文档顺序返回 body 元素序列：("h", level, title) / ("p", None, None) / ("tbl", None, None)。
    表格也算正文，避免"正文是一张表"的章节被判成空章节。"""
    seq = []
    try:
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        for child in d.element.body.iterchildren():
            tag = str(child.tag).rsplit("}", 1)[-1]
            if tag == "p":
                p = Paragraph(child, d)
                text = (p.text or "").strip()
                lvl = heading_level(p)
                if lvl is not None and text:
                    seq.append(("h", lvl, text))
                elif text:
                    seq.append(("p", None, None))
            elif tag == "tbl":
                tbl = Table(child, d)
                if any((c.text or "").strip() for r in tbl.rows for c in r.cells):
                    seq.append(("tbl", None, None))
    except Exception:
        for p in d.paragraphs:  # 降级：只看段落（老版本 python-docx）
            text = (p.text or "").strip()
            lvl = heading_level(p)
            if lvl is not None and text:
                seq.append(("h", lvl, text))
            elif text:
                seq.append(("p", None, None))
    return seq


def _check_docx(path, ext, checks):
    try:
        import docx  # python-docx
    except Exception:
        checks.append(_c("library_open", "skip", "python-docx unavailable"))
        return _result(path, ext, True, True, checks)
    try:
        d = docx.Document(path)
    except Exception as e:
        checks.append(_c("library_open", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    n = len(d.paragraphs)
    checks.append(_c("library_open", "pass", f"{n} paragraph(s)"))

    items, chars = _docx_texts(d)
    imgs, tbls = _docx_images(d), len(d.tables)
    if chars == 0 and imgs == 0 and tbls == 0:
        checks.append(_c("non_empty", "fail", "文档无文本、无图片、无表格"))
        return _result(path, ext, True, False, checks, paragraphs=n, text_chars=0, images=0)
    if chars == 0:
        # 与 pdf 分支同一口径（见 check_pdf 的 text_layer）：无文本 ≠ 空文档。
        # 图片型交付物（打印稿/扫描件/整页图表）天然没有文本层，判 fail 是误报 ——
        # 实测里这条误报把校验子代理直接引到"忽略技能结论、自己手写脚本"的路上。
        checks.append(_c("non_empty", "pass", f"无文本，但含 {imgs} 张图片 / {tbls} 张表格"))
        checks.append(_c("text_layer", "warn", "文档无可提取文本，文本内容无法核对（不做 OCR）"))
    else:
        checks.append(_c("non_empty", "pass", f"{chars} 个非空字符"))

    checks.append(_placeholder_check(items))

    # 空章节：某标题之后直到下一个同级或更高级标题之间没有任何正文。
    # 标题下紧跟更深一级标题（H1→H2）是正常结构，不报。
    # 必须按 body 的**文档顺序**遍历段落与表格：只看 d.paragraphs 会把"正文是一张表"
    # 的章节误判为空章节（真实误报来源）。
    empty = []
    seq = _docx_body_seq(d)
    heads = [(i, it[1], it[2]) for i, it in enumerate(seq) if it[0] == "h"]
    for idx, (si, lvl, title) in enumerate(heads):
        body = 0
        for kind, a, _b in seq[si + 1 :]:
            if kind == "h":
                if a <= lvl:
                    break
                continue  # 更深一级标题：本节有子结构，继续找正文
            body += 1
            break
        if body == 0:
            nxt = heads[idx + 1] if idx + 1 < len(heads) else None
            if nxt is None or nxt[1] <= lvl:
                empty.append(f"H{lvl} 「{title[:30]}」")
        if len(empty) >= MAX_HITS:
            break
    if empty:
        checks.append(_c("empty_section", "warn", "标题下无正文: " + "; ".join(empty)))
    else:
        checks.append(_c("empty_section", "pass"))

    # 页级取证：docx 不存页，页由静态页模型按 raw 下界口径推演 —— 下界只会漏报不会虚报。
    page_count = None
    try:
        from checkers import docx_pages

        pchecks, pinfo = docx_pages.checks_for(path, d)
        checks.extend(pchecks)
        page_count = pinfo.get("page_count_static")
    except Exception as e:
        checks.append(_c("blank_page", "skip", f"静态页模型不可用，空白页未校验: {str(e)[:160]}"))

    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks, paragraphs=n, text_chars=chars,
                   images=imgs, page_count=page_count)


def _shape_box(sh):
    try:
        if None in (sh.left, sh.top, sh.width, sh.height):
            return None
        return (int(sh.left), int(sh.top), int(sh.width), int(sh.height))
    except Exception:
        return None


def _check_pptx(path, ext, checks):
    try:
        from pptx import Presentation
    except Exception:
        checks.append(_c("library_open", "skip", "python-pptx unavailable"))
        return _result(path, ext, True, True, checks)
    try:
        prs = Presentation(path)
    except Exception as e:
        checks.append(_c("library_open", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    slides = list(prs.slides)
    n = len(slides)
    checks.append(_c("library_open", "pass", f"{n} slide(s)"))
    if n == 0:
        checks.append(_c("non_empty", "fail", "no slides"))
        return _result(path, ext, True, False, checks, slides=0, text_chars=0)

    sw, sh_ = int(prs.slide_width or 0), int(prs.slide_height or 0)
    texts, chars = [], 0
    empty_slides, oob, overlaps, tiny = [], [], [], []

    for si, slide in enumerate(slides, 1):
        shapes = list(iter_shapes(slide.shapes))
        has_visual, text_boxes = False, []
        for sh in shapes:
            t = shape_text(sh)
            if t:
                chars += len(t)
                texts.append((f"第{si}页", t))
            try:
                if sh.shape_type == 13 or sh.has_chart or sh.has_table:  # PICTURE / chart / table
                    has_visual = True
            except Exception:
                pass
            box = _shape_box(sh)
            if box and t:
                text_boxes.append((sh, box, t))
            # 出界：带文字的形状越界即确定性缺陷；纯装饰形状允许出血，超过 10% 才记 warn
            if box and sw and sh_:
                x, y, w, h = box
                over = max(-x, -y, x + w - sw, y + h - sh_)
                if over > 0:
                    ratio = over / max(sw, sh_)
                    if t and ratio > 0.02:
                        oob.append(f"第{si}页 文本框「{t[:20]}」越界 {ratio * 100:.0f}%")
                    elif not t and ratio > 0.10:
                        oob.append(f"第{si}页 装饰形状越界 {ratio * 100:.0f}%")
            # 字号过小
            try:
                if sh.has_text_frame:
                    for para in sh.text_frame.paragraphs:
                        for run in para.runs:
                            sz = run.font.size
                            if sz is not None and sz.pt < TINY_FONT_PT and (run.text or "").strip():
                                tiny.append(f"第{si}页 {sz.pt:.0f}pt「{run.text.strip()[:20]}」")
                                raise StopIteration
            except StopIteration:
                pass
            except Exception:
                pass
        # 文本框互相压盖：仅在两者都有文字、且相交超过较小者一半时记 warn
        for i in range(len(text_boxes)):
            for j in range(i + 1, len(text_boxes)):
                (_, (x1, y1, w1, h1), t1), (_, (x2, y2, w2, h2), t2) = text_boxes[i], text_boxes[j]
                iw = min(x1 + w1, x2 + w2) - max(x1, x2)
                ih = min(y1 + h1, y2 + h2) - max(y1, y2)
                if iw <= 0 or ih <= 0:
                    continue
                inter = iw * ih
                smaller = min(w1 * h1, w2 * h2) or 1
                if inter / smaller > 0.5:
                    overlaps.append(f"第{si}页「{t1[:15]}」与「{t2[:15]}」重叠 {inter / smaller * 100:.0f}%")
        if not any(shape_text(s) for s in shapes) and not has_visual:
            empty_slides.append(si)

    if empty_slides:
        checks.append(_c("empty_slide", "fail", "空白页（无文字也无图表）: " + ", ".join(map(str, empty_slides))))
    else:
        checks.append(_c("empty_slide", "pass"))
    checks.append(_placeholder_check(texts))
    if oob:
        checks.append(_c("out_of_bounds", "fail", "; ".join(oob[:MAX_HITS])))
    else:
        checks.append(_c("out_of_bounds", "pass"))
    if overlaps:
        checks.append(_c("text_overlap", "warn", "; ".join(overlaps[:MAX_HITS]) + "（可能为设计意图，需结合内容判断）"))
    else:
        checks.append(_c("text_overlap", "pass"))
    if tiny:
        checks.append(_c("tiny_font", "warn", f"字号小于 {TINY_FONT_PT:.0f}pt: " + "; ".join(tiny[:MAX_HITS])))
    else:
        checks.append(_c("tiny_font", "pass"))

    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks, slides=n, text_chars=chars)


def check_pdf(path, ext):
    checks = []
    reader_cls = None
    try:
        from pypdf import PdfReader as reader_cls  # type: ignore
    except Exception:
        try:
            from PyPDF2 import PdfReader as reader_cls  # type: ignore
        except Exception:
            reader_cls = None
    if reader_cls is None:
        try:
            with open(path, "rb") as f:
                head = f.read(5)
                f.seek(0)
                raw = f.read()
            ok = head == b"%PDF-"
            guess = len(re.findall(rb"/Type\s*/Page[^s]", raw))
            checks.append(_c("pdf_magic", "pass" if ok else "fail"))
            checks.append(
                _c("library_open", "skip", "pypdf/PyPDF2 unavailable，页数与文本未校验（/Type /Page 粗数 %d）" % guess)
            )
            return _result(path, ext, ok, ok, checks)
        except Exception as e:
            checks.append(_c("pdf_magic", "fail", str(e)))
            return _result(path, ext, False, False, checks)
    try:
        reader = reader_cls(path)
        pages = list(reader.pages)
    except Exception as e:
        checks.append(_c("library_open", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    checks.append(_c("library_open", "pass", f"{len(pages)} page(s)"))
    if not pages:
        checks.append(_c("non_empty", "fail", "0 pages"))
        return _result(path, ext, True, False, checks, page_count=0, text_chars=0)

    texts, chars, no_text, truly_empty = [], 0, [], []
    for i, pg in enumerate(pages, 1):
        try:
            t = (pg.extract_text() or "").strip()
        except Exception:
            t = ""
        chars += len(t)
        if t:
            texts.append((f"第{i}页", t))
            continue
        no_text.append(i)
        # 无文本不等于空白：矢量图/整页图片的页天然没有文本层。只有"既无文本、又无图片
        # 对象、内容流也几乎为空"才是确定性的空白页，避免把图表页误判成缺陷。
        blank = True
        try:
            res = pg.get("/Resources")
            res = res.get_object() if hasattr(res, "get_object") else res
            xo = (res or {}).get("/XObject")
            xo = xo.get_object() if hasattr(xo, "get_object") else xo
            if xo:
                blank = False
        except Exception:
            blank = False
        if blank:
            try:
                c = pg.get_contents()
                raw = c.get_data() if c is not None else b""
                if len(raw) >= 512:
                    blank = False
            except Exception:
                blank = False
        if blank:
            truly_empty.append(i)
    if truly_empty:
        checks.append(_c("non_empty", "fail", "空白页（无文本、无图形对象）: " + ", ".join(map(str, truly_empty[:MAX_HITS]))))
    elif chars == 0:
        checks.append(
            _c("text_layer", "warn", "全部页面均无可提取文本（整页图片/矢量输出），文本内容无法核对")
        )
        checks.append(_c("non_empty", "pass", "页面有图形内容"))
    else:
        checks.append(_c("non_empty", "pass", f"{chars} 个字符"))
        if no_text:
            checks.append(_c("text_layer", "warn", "无可提取文本的页: " + ", ".join(map(str, no_text[:MAX_HITS]))))
    checks.append(_placeholder_check(texts))

    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks, page_count=len(pages), text_chars=chars)


def check_html(path, ext):
    checks = []
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
    except Exception as e:
        checks.append(_c("read", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    if not content.strip():
        checks.append(_c("non_empty", "fail", "empty file"))
        return _result(path, ext, True, False, checks)
    checks.append(_c("non_empty", "pass", f"{len(content)} 字符（源码）"))
    low = content.lower()
    if "<html" in low or "<body" in low or "<!doctype" in low:
        checks.append(_c("html_structure", "pass"))
    else:
        checks.append(_c("html_structure", "warn", "no <html>/<body>/<!doctype>"))

    # 本地引用存在性：这里是文件系统层的权威结论。缺失即交付物不自洽 → fail。
    # 注意浏览器探针（probe）在 setContent 协议下对相对资源必然报失败，那是协议假象，
    # 判定本地引用是否缺失只认本项结论（见 SKILL.md 裁决规则）。
    missing = []
    base = os.path.dirname(os.path.abspath(path))
    for m in re.finditer(r'(?:src|href)\s*=\s*["\']([^"\']+)["\']', content, re.I):
        ref = m.group(1)
        if ref.startswith(("http://", "https://", "data:", "#", "mailto:", "//", "javascript:", "tel:")):
            continue
        if "{{" in ref or "${" in ref:  # 模板变量，运行时才成形
            continue
        rp = os.path.normpath(os.path.join(base, ref.split("?")[0].split("#")[0]))
        if not os.path.exists(rp):
            missing.append(ref)
        if len(missing) >= MAX_HITS:
            break
    if missing:
        checks.append(_c("local_refs", "fail", "引用的本地资源不存在: " + ", ".join(missing[:MAX_HITS])))
    else:
        checks.append(_c("local_refs", "pass"))

    # 源码层静态校验：内联 JS 语法（唯一 fail 级新增——一处失配整个 <script> 块不执行，
    # 该块所有图表全空）+ 图表容器/引库约定（warn，静态推断）+ 源码编码（warn）。
    # 渲染后的几何冲突（连线压穿卡片、浮框压住正文）不在这里：判据是渲染坐标，
    # 由 content_extract.py 的浏览器探针给出（见 SKILL.md 的 meta.probe.layout）。
    try:
        from checkers import html_static

        checks.extend(html_static.checks_for(path, content))
    except Exception as e:
        checks.append(_c("js_syntax", "skip", f"静态校验不可用，未覆盖: {str(e)[:160]}"))

    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks)


def check_text(path, ext):
    checks = []
    try:
        size = os.path.getsize(path)
    except Exception as e:
        checks.append(_c("exists", "fail", str(e)))
        return _result(path, ext, False, False, checks)
    if size == 0:
        checks.append(_c("non_empty", "fail", "0 bytes"))
        return _result(path, ext, True, False, checks)
    checks.append(_c("non_empty", "pass", f"{size} bytes"))
    if ext in ("md", "txt", "csv", "tsv", "json", "yaml", "yml"):
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = f.read()
            checks.append(_c("utf8", "pass"))
        except Exception as e:
            checks.append(_c("utf8", "warn", str(e)))
            data = None
        if ext == "json" and data is not None:
            try:
                json.loads(data)
                checks.append(_c("json_parse", "pass"))
            except Exception as e:
                checks.append(_c("json_parse", "fail", str(e)))
                return _result(path, ext, True, False, checks)
        if data:
            checks.append(_placeholder_check([(Path(path).name, data)]))
    else:
        # 不在本技能校验范围的类型（图片/音视频/压缩包/代码等）：只确认存在且非空，
        # 明确标 skip，避免"看起来合格"被当成通过。
        checks.append(_c("type_supported", "skip", f".{ext} 不在本技能校验范围，仅确认文件存在且非空"))
    fail = _has_fail(checks)
    return _result(path, ext, True, not fail, checks)


def run(path):
    if not os.path.exists(path):
        return _result(path, "", False, False, [_c("exists", "fail", "file not found")])
    ext = os.path.splitext(path)[1].lower().lstrip(".")
    if ext in ("pptx", "docx", "xlsx", "xlsm"):
        return check_office(path, ext)
    if ext == "pdf":
        return check_pdf(path, ext)
    if ext == "html":
        return check_html(path, ext)
    return check_text(path, ext)


def doctor(install=False):
    """依赖自检：一次输出全部依赖状态，替代调用方逐个探测。

    `install=True` 时用 deps.py 按**国内源优先**（清华 → 阿里 → 官方）把缺失项装进本技能
    缓存（`pip install --target`，不碰系统环境）：装不上再降级，且如实标注"未覆盖"。
    调用方不需要自己拼 pip 命令——默认官方源在国内经常装不上，结果就是维度静默丢失。
    """
    pkgs = (
        ("openpyxl", "openpyxl"),
        ("python-docx", "docx"),
        ("python-pptx", "pptx"),
        ("pypdf", "pypdf"),
        ("Pillow", "PIL"),
    )
    mods = {}
    for name, mod in pkgs:
        try:
            __import__(mod)
            mods[name] = "ok"
            continue
        except Exception:
            pass
        if install and deps is not None:
            ok, detail = deps.ensure_module(name, mod)
            mods[name] = ("ok（%s）" % detail) if ok else ("missing (%s)" % detail)
        else:
            mods[name] = "missing (加 --install 可按国内源优先自动安装)"

    bins = {}
    for b in ("soffice", "pdftoppm", "pdftotext", "dumate-browser-cli", "node"):
        p = shutil.which(b) or shutil.which(os.path.expanduser(f"~/.local/bin/{b}"))
        bins[b] = p or "missing"

    # html 内联 JS 语法自检需要一个真解析器：node 优先，否则 esprima（几百 KB，按国内源装）
    parser, mirrors = {"kind": None, "detail": "deps.py 不可用"}, []
    if deps is not None:
        try:
            parser = deps.ensure_js_parser(install=True) if install else deps.probe_js_parser()
            mirrors = deps.mirror_order()
        except Exception as e:
            parser = {"kind": None, "detail": f"探测异常: {str(e)[:120]}"}

    layout_on = os.environ.get("DUMATE_DELIVERABLE_LAYOUT_AUDIT", "").strip() not in ("0", "false")
    caps = {
        "structure_check": "ok",
        "content_extract": "ok",
        "link_probe": "ok（需外网；无外网时整批记 unknown）",
        "xlsx_formula_recalc": "ok" if bins["soffice"] != "missing" else "degraded (无 soffice：仅能读缓存值)",
        "docx_page_model": "ok（静态页模型：空白页 / 单段超页；页码为推演值，可能与实际渲染差 ±1）",
        "html_render": (
            "ok（extract 抽取 html 时驱动无头浏览器，产出渲染后文本与运行时探针）"
            if bins["dumate-browser-cli"] != "missing"
            else "unavailable (无 dumate-browser-cli：html 退化为源码去标签，运行时缺陷未覆盖)"
        ),
        "html_js_syntax": (
            f"ok（{parser['kind']}：{parser.get('detail', '')}）"
            if parser.get("kind")
            else f"degraded (无 node/esprima：只能做启发式括号配平，语法按未覆盖处理 —— {parser.get('detail', '')})"
        ),
        "html_layout_audit": (
            "ok（渲染后几何审计：装饰线与内容框、浮框与正文行盒，随 html 渲染一并产出 probe.layout）"
            if bins["dumate-browser-cli"] != "missing" and layout_on
            else (
                "disabled (DUMATE_DELIVERABLE_LAYOUT_AUDIT=0)"
                if not layout_on
                else "unavailable (无 dumate-browser-cli：渲染后几何未覆盖)"
            )
        ),
        "visual_audit": "disabled (本流程不做多模态视觉排版审核)",
    }
    return {
        "python": sys.version.split()[0],
        "modules": mods,
        "binaries": bins,
        "capabilities": caps,
        "install_mirrors": mirrors,
    }


def main():
    ap = argparse.ArgumentParser(description="Deliverable structure/content check (read-only).")
    ap.add_argument("file", nargs="?")
    ap.add_argument("--doctor", action="store_true", help="依赖自检并退出")
    ap.add_argument("--install", action="store_true",
                    help="与 --doctor 合用：缺失依赖按国内源优先自动安装到本技能缓存")
    args = ap.parse_args()
    if args.doctor or args.install:
        try:
            print(json.dumps(doctor(install=args.install), ensure_ascii=False))
        except Exception as e:  # doctor 也不许抛栈
            print(json.dumps({"error": f"doctor failed: {str(e)[:200]}"}, ensure_ascii=False))
        return
    if not args.file:
        print(json.dumps(
            {"error": "usage: defect_check.py <file> | defect_check.py --doctor [--install]"},
            ensure_ascii=False,
        ))
        sys.exit(2)
    try:
        result = run(args.file)
    except Exception as e:
        result = _result(args.file, "", False, False, [_c("unexpected", "fail", str(e))])
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
